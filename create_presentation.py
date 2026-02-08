from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt(filename="NK_Automobiles_Final_Presentation.pptx"):
    prs = Presentation()

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[0] # 0 is Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "NK Automobiles"
    subtitle.text = "Comprehensive E-Commerce Platform for Bike Spare Parts\nProject Presentation"

    # --- Introduction ---
    slide_layout = prs.slide_layouts[1] # 1 is Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Introduction"
    content.text = ("NK Automobiles is a dedicated e-commerce web application designed to simplify "
                    "the process of purchasing genuine two-wheeler spare parts online.\n\n"
                    "It bridges the gap between bike owners and authentic spare part suppliers, "
                    "providing a seamless shopping experience with secure payment options and "
                    "reliable delivery tracking.")

    # --- Problem Statement ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Problem Statement"
    content.text = ("Current Challenges in the Market:\n"
                    "- Difficulty in finding genuine spare parts for specific bike models.\n"
                    "- Uncertainty about product quality and authenticity.\n"
                    "- Lack of transparent pricing and inventory information.\n"
                    "- Inconvenient offline purchasing process requiring physical store visits.")

    # --- Solution & Key Features (User) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Solution: Key User Features"
    content.text = ("- Extensive Product Catalog: Categorized by brand and part type.\n"
                    "- Smart Search & Filters: Easily find parts by bike model or price range.\n"
                    "- Secure Authentication: User registration, login, and profile management.\n"
                    "- Interactive Shopping: Wishlist, Shopping Cart, and Product Reviews.\n"
                    "- Multiple Payment Options: Cash on Delivery (COD) and UPI QR Code payments.")

    # --- Solution & Key Features (Admin) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Solution: Admin Dashboard"
    content.text = ("- Product Management: Add, edit, delete products, and manage stock levels.\n"
                    "- Order Management: View order details, update status (Pending, Shipped, Delivered).\n"
                    "- User Management: View registered users and their activity.\n"
                    "- Coupon Management: Create and manage discount codes.\n"
                    "- Analytics: Dashboard overview of total sales, orders, and user growth.")

    # --- Technology Stack ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Technology Stack"
    content.text = ("Frontend:\n"
                    "- HTML5, CSS3, Bootstrap 5 (Responsive Design)\n"
                    "- JavaScript (Interactive elements)\n\n"
                    "Backend:\n"
                    "- Python (Core Logic)\n"
                    "- Flask (Web Framework)\n\n"
                    "Database:\n"
                    "- MySQL (Data Storage & Management)")

    # --- Payment Workflow (Special Focus) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Secure Payment Workflow"
    content.text = ("Hybrid Payment System:\n"
                    "1. UPI QR Code Integration:\n"
                    "   - Users scan a dynamic QR code generated for the order amount.\n"
                    "   - Mandatory Transaction ID verification to prevent fraud.\n"
                    "   - Duplicate transaction check to ensure unique payments.\n\n"
                    "2. Cash on Delivery (COD):\n"
                    "   - Traditional option for user convenience.\n"
                    "   - Verified via OTP/Phone confirmation (process feature).")

    # --- Database Schema (ER Diagram) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Database Schema (ER Relationships)"
    content.text = ("1. Users Entity:\n"
                    "   - One-to-Many relationship with Orders and Reviews.\n"
                    "   - One-to-One relationship with Cart.\n\n"
                    "2. Products Entity:\n"
                    "   - Belongs to one Category.\n"
                    "   - Has many Reviews and Order Items.\n\n"
                    "3. Orders Entity:\n"
                    "   - Linked to one User.\n"
                    "   - Contains multiple Order Items (Many-to-Many via Order_Items table).\n\n"
                    "4. Core Tables: Users, Products, Categories, Orders, Cart, Reviews.")

    # --- User Workflow (Flowchart) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "User Workflow (Flowchart)"
    content.text = ("1. Start: User visits Home Page.\n"
                    "2. Browse: User searches for products or filters by category.\n"
                    "3. Select: User views Product Details and clicks 'Add to Cart'.\n"
                    "4. Authenticate: User Logs in or Registers (if not active).\n"
                    "5. Checkout: User enters shipping address and chooses payment.\n"
                    "6. Payment: \n"
                    "   - If UPI: Scan QR -> Pay -> Enter Transaction ID.\n"
                    "   - If COD: Confirm Order directly.\n"
                    "7. End: Order Confirmed -> Email Sent -> Order Tracking.")

    # --- Future Enhancements ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Future Enhancements"
    content.text = ("- Automated Payment Gateway Integration (Razorpay/Stripe).\n"
                    "- AI-Powered Recommendation System for related parts.\n"
                    "- Mobile Application for Android/iOS.\n"
                    "- Live Chat Support integration.\n"
                    "- Multi-language support for broader accessibility.")

    # --- Conclusion ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Conclusion"
    content.text = ("NK Automobiles successfully addresses the market need for a reliable online spare parts store. "
                    "By combining a user-friendly interface with robust backend management, it ensures efficient operations "
                    "and customer satisfaction.")

    # --- Thank You ---
    slide_layout = prs.slide_layouts[0] # Title Slide layout again for ending
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Thank You"
    subtitle.text = "Questions & Answers"

    prs.save(filename)
    print(f"Presentation saved as {filename}")

if __name__ == "__main__":
    create_ppt()
